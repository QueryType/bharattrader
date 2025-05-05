"""
Document Processor Module

This module handles the conversion of various file formats to markdown text.
Supported formats: txt, pdf, docx, pptx, xlsx, images
"""

import os
import re
import logging
import base64
import json
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any
import datetime

# Load environment variables from .env file
try:
    from dotenv import load_dotenv
    load_dotenv()  # Load variables from .env
    ENV_LOADED = True
except ImportError:
    ENV_LOADED = False
    logging.warning("dotenv not found, environment variables must be set manually")

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Get model IDs and config from environment variables
OPENAI_VISION_MODEL = os.environ.get("OPENAI_VISION_MODEL", "gpt-4-vision-preview")
OPENAI_TEXT_MODEL = os.environ.get("OPENAI_TEXT_MODEL", "gpt-4-turbo")
# Flag to enable/disable LLM prompt logging (default: enabled)
ENABLE_LOGGING = os.environ.get("ENABLE_LLM_LOGGING", "true").lower() == "true"

# Try to import optional dependencies, with graceful fallbacks
try:
    import fitz  # PyMuPDF
    PDF_EXTRACTOR = "pymupdf"
except ImportError:
    PDF_EXTRACTOR = None
    logger.warning("PyMuPDF not found. PDF extraction will be limited.")

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not found. DOCX extraction will be unavailable.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not found. PPTX extraction will be unavailable.")

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False
    logger.warning("pandas not found. XLSX extraction will be unavailable.")

try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    logger.warning("PIL or pytesseract not found. OCR will be unavailable.")

try:
    import openai
    OPENAI_AVAILABLE = True
except ImportError:
    OPENAI_AVAILABLE = False
    logger.warning("OpenAI library not found. Advanced image analysis will be unavailable.")


def extract_from_txt(file_path: str) -> str:
    """Extract text from a plain text file.

    Args:
        file_path: Path to the text file

    Returns:
        Extracted text content
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # Try with different encodings if utf-8 fails
        try:
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {str(e)}")
            return f"ERROR: Could not read {file_path} due to encoding issues."


def extract_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file.

    Args:
        file_path: Path to the PDF file

    Returns:
        Extracted text content
    """
    if PDF_EXTRACTOR == "pymupdf":
        try:
            text_content = []
            with fitz.open(file_path) as doc:
                for page_num, page in enumerate(doc):
                    text = page.get_text()
                    text_content.append(f"# Page {page_num + 1}\n\n{text}\n\n")
            return "\n".join(text_content)
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {str(e)}")
            return f"ERROR: Could not extract text from {file_path}."
    else:
        logger.error("No PDF extraction library available")
        return "ERROR: PDF extraction requires PyMuPDF. Please install with: pip install pymupdf"


def extract_from_docx(file_path: str) -> str:
    """Extract text from a DOCX file.

    Args:
        file_path: Path to the DOCX file

    Returns:
        Extracted text content
    """
    if not DOCX_AVAILABLE:
        return "ERROR: DOCX extraction requires python-docx. Please install with: pip install python-docx"
    
    try:
        doc = docx.Document(file_path)
        full_text = []
        
        for para in doc.paragraphs:
            full_text.append(para.text)
            
        # Add tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text for cell in row.cells])
                full_text.append(f"| {row_text} |")
        
        return "\n\n".join(full_text)
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {file_path}: {str(e)}")
        return f"ERROR: Could not extract text from {file_path}."


def extract_from_pptx(file_path: str) -> str:
    """Extract text from a PPTX file.

    Args:
        file_path: Path to the PPTX file

    Returns:
        Extracted text content
    """
    if not PPTX_AVAILABLE:
        return "ERROR: PPTX extraction requires python-pptx. Please install with: pip install python-pptx"
    
    try:
        presentation = Presentation(file_path)
        text_content = []
        
        for slide_num, slide in enumerate(presentation.slides):
            slide_text = []
            slide_text.append(f"# Slide {slide_num + 1}")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            text_content.append("\n\n".join(slide_text))
        
        return "\n\n---\n\n".join(text_content)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {file_path}: {str(e)}")
        return f"ERROR: Could not extract text from {file_path}."


def extract_from_xlsx(file_path: str) -> str:
    """Extract data from an Excel file.

    Args:
        file_path: Path to the Excel file

    Returns:
        Extracted data as markdown tables
    """
    if not PANDAS_AVAILABLE:
        return "ERROR: Excel extraction requires pandas. Please install with: pip install pandas openpyxl"
    
    try:
        result = []
        # Read all sheets
        excel_file = pd.ExcelFile(file_path)
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            # Convert to markdown table
            md_table = f"## Sheet: {sheet_name}\n\n"
            md_table += df.to_markdown(index=False)
            result.append(md_table)
        
        return "\n\n---\n\n".join(result)
    except Exception as e:
        logger.error(f"Error extracting data from Excel {file_path}: {str(e)}")
        return f"ERROR: Could not extract data from {file_path}."


def log_llm_prompt(
    company_name: str, 
    phase: str, 
    section: str, 
    messages: List[Dict[str, Any]],
    model: str,
    temperature: float,
    max_tokens: int,
    run_timestamp: Optional[str] = None
) -> None:
    """Log the prompt sent to the LLM.
    
    Args:
        company_name: Name of the company
        phase: Phase of processing (e.g., 'document_processing')
        section: Section being generated (e.g., 'image_analysis')
        messages: Messages sent to the LLM
        model: Model name
        temperature: Temperature setting
        max_tokens: Max tokens setting
        run_timestamp: Optional timestamp to use for the log filename. If provided, 
                       appends to an existing log file with this timestamp.
    """
    if not ENABLE_LOGGING:
        logger.info("LLM logging is disabled. Skipping log entry.")
        return
    
    # Create logs directory
    company_logs_dir = Path(f"company_data/{company_name}/logs")
    company_logs_dir.mkdir(exist_ok=True, parents=True)
    
    # Generate timestamp for the log file or use provided one
    timestamp = run_timestamp if run_timestamp else datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    log_filename = f"{company_name}_{phase}_{timestamp}.log"
    log_path = company_logs_dir / log_filename
    
    # Prepare log entry
    log_entry = {
        "timestamp": datetime.datetime.now().isoformat(),
        "company": company_name,
        "phase": phase,
        "section": section,
        "model": model,
        "temperature": temperature,
        "max_tokens": max_tokens,
        "messages": messages
    }
    
    # Append to log file
    try:
        # Create file if it doesn't exist
        if not log_path.exists():
            with open(log_path, 'w', encoding='utf-8') as f:
                f.write(f"# LLM Interaction Log for {company_name}\n")
                f.write(f"# Phase: {phase}\n")
                f.write(f"# Created: {timestamp}\n\n")
        
        # Append log entry
        with open(log_path, 'a', encoding='utf-8') as f:
            f.write(f"\n## {section} - {datetime.datetime.now().isoformat()}\n")
            f.write(json.dumps(log_entry, indent=2))
            f.write("\n\n---\n\n")
        
        logger.info(f"Logged LLM prompt for {company_name}/{phase}/{section} to {log_path}")
    except Exception as e:
        logger.error(f"Failed to log LLM prompt: {str(e)}")


def extract_from_image(file_path: str, run_timestamp: Optional[str] = None) -> str:
    """Extract text from an image using OCR.

    Args:
        file_path: Path to the image file
        run_timestamp: Optional timestamp for consistent log file naming

    Returns:
        Extracted text content
    """
    if not OCR_AVAILABLE:
        return "ERROR: Image extraction requires Pillow and pytesseract. Please install with: pip install Pillow pytesseract"
    
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        
        # Use OpenAI for better image understanding if available
        if OPENAI_AVAILABLE and os.environ.get("OPENAI_API_KEY"):
            try:
                # Get company name from file path
                file_path_obj = Path(file_path)
                company_name = file_path_obj.parent.name
                
                # Create a timestamp for this processing run if not provided
                if run_timestamp is None:
                    run_timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
                
                # Try to get a better description using OpenAI's vision capabilities
                client = openai.Client(api_key=os.environ["OPENAI_API_KEY"])
                with open(file_path, "rb") as image_file:
                    base_image = image_file.read()
                
                # Prepare messages
                messages = [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "Describe this image in detail, focusing on any financial data, charts, or business information visible."},
                            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64.b64encode(base_image).decode('utf-8')}"}}
                        ]
                    }
                ]
                
                # Log the prompt
                log_llm_prompt(
                    company_name=company_name,
                    phase="document_processing",
                    section=f"image_analysis_{file_path_obj.stem}",
                    messages=messages,
                    model=OPENAI_VISION_MODEL,  # Use global variable
                    temperature=0.3,
                    max_tokens=300,
                    run_timestamp=run_timestamp
                )
                
                response = client.chat.completions.create(
                    model=OPENAI_VISION_MODEL,  # Use global variable
                    messages=messages,
                    temperature=0.3,
                    max_tokens=300
                )
                vision_description = response.choices[0].message.content
                return f"## OCR Text:\n\n{text}\n\n## Image Analysis:\n\n{vision_description}"
            except Exception as e:
                logger.warning(f"OpenAI vision processing failed: {str(e)}")
                return f"## OCR Text:\n\n{text}"
        else:
            return f"## OCR Text:\n\n{text}"
    except Exception as e:
        logger.error(f"Error extracting text from image {file_path}: {str(e)}")
        return f"ERROR: Could not extract text from {file_path}."


def convert_to_markdown(file_path: str, run_timestamp: Optional[str] = None) -> Tuple[str, str]:
    """Convert various file formats to markdown text.

    Args:
        file_path: Path to the file
        run_timestamp: Optional timestamp for consistent log naming across a run

    Returns:
        Tuple of (markdown_content, file_name)
    """
    file_path = Path(file_path)
    file_ext = file_path.suffix.lower()[1:]  # Remove the dot
    file_name = file_path.stem
    
    content = f"# {file_name}\n\n"
    content += f"Source: {file_path}\n"
    content += f"Processed on: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n---\n\n"
    
    try:
        if file_ext == "txt":
            extracted = extract_from_txt(str(file_path))
        elif file_ext == "pdf":
            extracted = extract_from_pdf(str(file_path))
        elif file_ext == "docx":
            extracted = extract_from_docx(str(file_path))
        elif file_ext == "pptx":
            extracted = extract_from_pptx(str(file_path))
        elif file_ext in ["xlsx", "xls"]:
            extracted = extract_from_xlsx(str(file_path))
        elif file_ext in ["jpg", "jpeg", "png", "gif", "bmp"]:
            extracted = extract_from_image(str(file_path), run_timestamp)
        else:
            extracted = f"Unsupported file format: {file_ext}"
            logger.warning(f"Unsupported file format: {file_ext}")
        
        content += extracted
        
    except Exception as e:
        logger.error(f"Error processing {file_path}: {str(e)}")
        content += f"ERROR: Failed to process file {file_path}. Exception: {str(e)}"
    
    return content, f"{file_name}.md"


def process_company_folder(company_folder: str) -> List[Tuple[str, str]]:
    """Process all files in a company folder.

    Args:
        company_folder: Path to the company folder

    Returns:
        List of tuples (markdown_content, markdown_file_path)
    """
    logger.info(f"Processing company folder: {company_folder}")
    company_path = Path(company_folder)
    
    if not company_path.exists() or not company_path.is_dir():
        logger.error(f"Company folder does not exist: {company_folder}")
        return []
    
    # Get company name from folder name
    company_name = company_path.name
    
    # Create output folders proactively
    output_folder = company_path / "processed"
    output_folder.mkdir(exist_ok=True)
    
    # Create logs directory proactively
    logs_folder = company_path / "logs"
    logs_folder.mkdir(exist_ok=True)
    logger.info(f"Ensured logs directory exists: {logs_folder}")
    
    # Create a single timestamp for this processing run
    run_timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    
    results = []
    
    # Process all files in the folder
    for file_path in company_path.glob("*"):
        if file_path.is_file() and not file_path.name.startswith('.') and not file_path.name.endswith('.md'):
            logger.info(f"Processing file: {file_path}")
            
            # Convert the file to markdown using the common run timestamp
            markdown_content, markdown_name = convert_to_markdown(str(file_path), run_timestamp)
            
            # Save the markdown file
            markdown_path = output_folder / markdown_name
            with open(markdown_path, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            results.append((markdown_content, str(markdown_path)))
    
    logger.info(f"Processed {len(results)} files for company: {company_name}")
    return results